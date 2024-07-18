from docx import Document
from faker import Faker
import random
import pyautogui
from pptx import Presentation
from pptx.util import Inches
import logging
from pathlib import Path
from utils import util

random.seed(43)

logger = util.create_logger(Path(__file__).name)

fake = Faker()


def create_doc():
    doc = Document()
    # fack name
    name = fake.name()
    address = fake.address()
    doc.add_heading(f"{name}: {address}")

    times = random.random(5, 20)
    for _ in range(times):
        # fake context
        text = fake.text()
        sentence = fake.sentence()
        doc.add_paragraph(sentence, styple="ListBullet")
        doc.add_paragraph(text)
    logger.info("create a document file: {name}.docx")
    doc.save("{name}.docx")

def modify_doc():
    docx_files = list(Path.cwd().glob("*.docx"))
    if docx_files:
        docx_path = random.choice(docx_files)
        doc = Document(docx_path)
        first_para = doc.paragraphs[0]
        first_para.text = fake.sentence()
        doc.add_paragraph(fake.sentence())
        logger.info("modified the powerpoint file: {pptx_path.name}.docx")
        doc.save(docx_path)


def delete_doc():
    docx_files = list(Path.cwd().glob("*.docx"))
    if docx_files:
        docx_path = random.choice(docx_files)
        logger.info("deleted the powerpoint file: {docx_path.name}.pptx")
        docx_path.unlink()

def create_ppt():
    # create a ppt object
    prs = Presentation()
    # add a title slide --- 0 is the index for title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = fake.text()
    subtitle.text = fake.text()

    # add a slide with a title and content layout
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = fake.text()
    content.text = fake.sentence()

    # add another slide with a different layout
    page_index = random.randint(1,100)
    slide_layout = prs.slide_layouts[page_index] 
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = fake.text()

    # add an image to the slide
    img_path = Path.cwd().joinpath("tree.jpeg").as_posix()
    name = fake.text()
    save_path = Path.cwd().joinpath(f"{name}.pptx").as_posix()
    left = Inches(1)
    top = Inches(1)
    height = Inches(2)
    slide.shape.add_picture(img_path, left, top, height=height)
    logger.info("create a powerpoint file: {name}.pptx")
    prs.save(save_path)


def modify_ppt():
    pptx_files = list(Path.cwd().glob("*.pptx"))
    if pptx_files:
        pptx_path = random.choice(pptx_files)
        prs = Presentation(pptx_path)

        # make changes to title and context
        first_slide = prs.slides[0]
        title = first_slide.shapes.title
        title.text = fake.text()
        logger.info("modified the powerpoint file: {pptx_path.name}.pptx")
        prs.save(pptx_path)

def delete_ppt():
    pptx_files = list(Path.cwd().glob("*.pptx"))
    if pptx_files:
        pptx_path = random.choice(pptx_files)
        logger.info("deleted the powerpoint file: {pptx_path.name}.pptx")
        pptx_path.unlink()


def automate_gui():
    pyautogui.moveTo(100, 100, duration=1)
    pyautogui.click()
    pyautogui.write("Automating office work with Python", interval=0.1)
    pyautogui.press("enter")
