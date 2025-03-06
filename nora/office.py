import sys
from pathlib import Path
sys.path.insert(0, Path(sys.path[0]).parent.as_posix())
import random
import os

try:
    from faker import Faker
    from pptx import Presentation
    from docx import Document
    from pptx.util import Inches
    from openpyxl import Workbook, load_workbook
    import pyautogui
    
except:
    os.system("pip3 install faker==26")
    os.system("pip3 install python-pptx==0.6.23")
    os.system("pip3 install python-docx==1.1.2")
    os.system("pip3 install openpyxl==3.1.5")
    os.system("pip3 install pyautogui==0.9.54")


from pathlib import Path
from utils import util
import string


random.seed(43)

logger = util.create_logger(Path.cwd().parent.joinpath("logs", Path(__file__).name))

fake = Faker("en_US")

def create_doc():
    doc = Document()
    # fack name
    name = fake.name()
    address = fake.address()
    doc.add_heading(f"{name}: {address}")

    times = random.random(5, 20)
    for _ in range(times):
        # fake context
        text = fake.text()
        sentence = fake.sentence()
        doc.add_paragraph(sentence, styple="ListBullet")
        doc.add_paragraph(text)
    logger.info("create a document file: {name}.docx")
    doc.save("{name}.docx")


def modify_doc():
    docx_files = list(Path.cwd().glob("*.docx"))
    if docx_files:
        # randomly modify a document
        docx_path = random.choice(docx_files)
        doc = Document(docx_path)
        first_para = doc.paragraphs[0]
        first_para.text = fake.sentence()
        doc.add_paragraph(fake.sentence())
        logger.info("modified the powerpoint file: {pptx_path.name}.docx")
        doc.save(docx_path)


def delete_doc():
    docx_files = list(Path.cwd().glob("*.docx"))
    if docx_files:
        # randomly pick one file to delete
        docx_path = random.choice(docx_files)
        logger.info("deleted the powerpoint file: {docx_path.name}.pptx")
        docx_path.unlink()


def create_ppt():
    # create a ppt object
    prs = Presentation()
    # add a title slide --- 0 is the index for title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = fake.text()
    subtitle.text = fake.text()

    # add a slide with a title and content layout
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = fake.text()
    content.text = fake.sentence()

    # add another slide with a different layout
    page_index = random.randint(1,100)
    slide_layout = prs.slide_layouts[page_index] 
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = fake.text()

    # add an image to the slide
    img_path = Path.cwd().joinpath("tree.jpeg").as_posix()
    name = fake.text()
    save_path = Path.cwd().joinpath(f"{name}.pptx").as_posix()
    left = Inches(1)
    top = Inches(1)
    height = Inches(2)
    slide.shape.add_picture(img_path, left, top, height=height)
    logger.info("create a powerpoint file: {name}.pptx")
    prs.save(save_path)


def modify_ppt():
    pptx_files = list(Path.cwd().glob("*.pptx"))
    if pptx_files:
        pptx_path = random.choice(pptx_files)
        prs = Presentation(pptx_path)

        # make changes to title and context
        first_slide = prs.slides[0]
        title = first_slide.shapes.title
        title.text = fake.text()
        logger.info("modified the powerpoint file: {pptx_path.name}.pptx")
        prs.save(pptx_path)


def delete_ppt():
    pptx_files = list(Path.cwd().glob("*.pptx"))
    if pptx_files:
        # randomly pick one file to delete
        pptx_path = random.choice(pptx_files)
        logger.info("deleted the powerpoint file: {pptx_path.name}.pptx")
        pptx_path.unlink()


def create_xls():
    # create a new workbook
    wb = Workbook()
    ws = wb.active
    
    # define the column headers
    column_len = random.randint(5,20)
    columns = [fake.name() for i in range(column_len)]

    # write the headers to first rows
    for i, col in enumerate(columns, start=1):
        ws.cell(row=1, column=i, value=col)    
    
    # define random data for each column
    row_len = random.randint(5,30)
    for row in range(2, row_len + 2):
        selected_columns = random.sample(range(1, column_len + 1), \
                                         k=random.randint(1, column_len))
    
        # Write random data to the selected columns
        for col_index in selected_columns:
            random_data = fake.text(max_nb_chars=20)  # Generate random text with a max of 20 characters
            ws.cell(row=row, column=col_index, value=random_data)

    # Save the workbook to a file
    file_name = fake.name()
    wb.save("{file_name}.xlsx")


def modify_xls():
    xlsx_files = list(Path.cwd().glob("*.xlsx"))
    if xlsx_files:
        # randomly modify a document
        xlsx_path = random.choice(xlsx_files)
        wb = load_workbook(xlsx_path)
        ws = wb.active  

        random_column = random.choice(string.ascii_uppercase)
        random_row = random.randint(5,100)
        ws[f'{random_column}{random_row}'] = fake.text()
        file_name = xlsx_path.stem
        wb.save("{file_name}.xlsx")

def delete_xls():
    xlsx_files = list(Path.cwd().glob("*.xlsx"))
    if xlsx_files:
        # randomly pick one file to delete
        xlsx_path = random.choice(xlsx_files)
        logger.info("deleted the excel file: {xlsx_path.sten}.xlsx")
        xlsx_path.unlink()


def automate_gui():
    pyautogui.moveTo(100, 100, duration=1)
    pyautogui.click()
    pyautogui.write("Automating office work with Python", interval=0.1)
    pyautogui.press("enter")


if __name__ == "__main__":
    # document manipulations
    create_doc()
    modify_doc()
    delete_doc()

    # powerpoint manipulation
    create_ppt()
    modify_ppt()
    delete_ppt()

    # excel manipulation
    create_xls()
    modify_xls()
    delete_xls()

    # click event
    automate_gui()


